from datetime import datetime
from pathlib import Path
from urllib.request import Request, urlopen

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WIDTH = 13.333
HEIGHT = 7.5

# shadcn-inspired tokens
BACKGROUND = RGBColor(250, 250, 250)
BACKGROUND_SOFT = RGBColor(245, 245, 245)
CARD_BG = RGBColor(255, 255, 255)
BORDER = RGBColor(228, 228, 231)
TEXT = RGBColor(24, 24, 27)
MUTED = RGBColor(113, 113, 122)
PRIMARY = RGBColor(24, 24, 27)
PRIMARY_SOFT = RGBColor(39, 39, 42)
PRIMARY_FG = RGBColor(250, 250, 250)
ACCENT = RGBColor(37, 99, 235)
ACCENT_SOFT = RGBColor(239, 246, 255)
SUCCESS = RGBColor(22, 163, 74)
WARNING = RGBColor(217, 119, 6)
WHITE = RGBColor(255, 255, 255)

# Keep these aliases to minimize per-slide changes
BLUE = PRIMARY_SOFT
BLUE_DARK = PRIMARY
BLUE_LIGHT = RGBColor(241, 245, 249)
GREEN = SUCCESS
ORANGE = WARNING

FONT = "Inter"
FONT_HEAD = "Inter SemiBold"
FONT_SCRIPT = "Segoe Script"

ASSETS = Path("assets")
RAW = ASSETS / "raw"
PROCESSED = ASSETS / "processed"

IMAGE_URLS = {
    "cover": "https://loremflickr.com/1920/1080/students,classroom,africa?lock=101",
    "problem": "https://loremflickr.com/1920/1080/teacher,stress?lock=102",
    "mission": "https://loremflickr.com/1920/1080/teacher,laptop?lock=103",
    "market": "https://loremflickr.com/1920/1080/africa,city,education?lock=104",
    "value": "https://loremflickr.com/1920/1080/teacher,students?lock=105",
    "traction": "https://loremflickr.com/1920/1080/students,computer,lab?lock=106",
    "technology": "https://loremflickr.com/1920/1080/school,technology,students?lock=107",
    "closing": "https://loremflickr.com/1920/1080/africa,education,map?lock=108",
    "team_1": "https://randomuser.me/api/portraits/men/32.jpg",
    "team_2": "https://randomuser.me/api/portraits/men/75.jpg",
    "team_3": "https://randomuser.me/api/portraits/men/41.jpg",
    "team_4": "https://randomuser.me/api/portraits/women/44.jpg",
}


def safe_name(key):
    return "".join(ch if ch.isalnum() or ch in ("_", "-") else "_" for ch in key)


def fetch_images():
    RAW.mkdir(parents=True, exist_ok=True)
    for key, url in IMAGE_URLS.items():
        dest = RAW / f"{safe_name(key)}.jpg"
        if dest.exists() and dest.stat().st_size > 0:
            continue
        req = Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urlopen(req, timeout=30) as r:
            data = r.read()
        dest.write_bytes(data)


def cropped_image(key, width_in, height_in):
    PROCESSED.mkdir(parents=True, exist_ok=True)
    src = RAW / f"{safe_name(key)}.jpg"
    out = PROCESSED / f"{safe_name(key)}_{int(width_in * 100)}x{int(height_in * 100)}.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return out

    ratio = width_in / height_in
    with Image.open(src) as img:
        img = img.convert("RGB")
        src_ratio = img.width / img.height
        if src_ratio > ratio:
            new_w = int(img.height * ratio)
            left = (img.width - new_w) // 2
            img = img.crop((left, 0, left + new_w, img.height))
        else:
            new_h = int(img.width / ratio)
            top = (img.height - new_h) // 2
            img = img.crop((0, top, img.width, top + new_h))
        px_w = max(900, int(width_in * 220))
        px_h = max(500, int(height_in * 220))
        img = img.resize((px_w, px_h), Image.Resampling.LANCZOS)
        img.save(out, format="JPEG", quality=92)
    return out


def style_text(shape, text, size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT, italic=False, font=FONT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font


def style_bullets(shape, items, size=20, color=TEXT, spacing=7):
    tf = shape.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing)


def add_background(slide, soft=False):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(WIDTH), Inches(HEIGHT))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BACKGROUND_SOFT if soft else BACKGROUND
    bg.line.fill.background()

    corner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.3), Inches(5.4), Inches(3.4))
    corner.fill.solid()
    corner.fill.fore_color.rgb = BLUE_LIGHT
    corner.fill.transparency = 0.5
    corner.line.fill.background()

    corner2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.7), Inches(0.2), Inches(4.4), Inches(2.8))
    corner2.fill.solid()
    corner2.fill.fore_color.rgb = RGBColor(244, 244, 245)
    corner2.fill.transparency = 0.4
    corner2.line.fill.background()


def add_header(slide, title, right_text=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(WIDTH), Inches(0.92))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.82), Inches(WIDTH), Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.34), Inches(0.13), Inches(8.7), Inches(0.58))
    style_text(t, title, size=28, bold=True, color=PRIMARY_FG, font=FONT_HEAD)

    if right_text:
        rt = slide.shapes.add_textbox(Inches(8.5), Inches(0.18), Inches(4.45), Inches(0.5))
        style_text(rt, right_text, size=20, bold=True, color=PRIMARY_FG, align=PP_ALIGN.RIGHT, italic=True, font=FONT_HEAD)


def add_tagline(slide, text, color=BLUE_DARK):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.95), Inches(WIDTH), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.3), Inches(7.04), Inches(12.7), Inches(0.35))
    style_text(tx, text, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True, font=FONT_HEAD)


def add_photo(slide, key, left, top, width, height, rounded=True):
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = CARD_BG
    frame.line.color.rgb = BORDER

    img_path = cropped_image(key, width - 0.12, height - 0.12)
    slide.shapes.add_picture(
        str(img_path),
        Inches(left + 0.06),
        Inches(top + 0.06),
        width=Inches(width - 0.12),
        height=Inches(height - 0.12),
    )


def add_full_photo(slide, key):
    image = cropped_image(key, WIDTH, HEIGHT)
    slide.shapes.add_picture(str(image), 0, 0, width=Inches(WIDTH), height=Inches(HEIGHT))


def card(slide, left, top, width, height, color=CARD_BG):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top + 0.03), Inches(width), Inches(height))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(212, 212, 216)
    shadow.fill.transparency = 0.82
    shadow.line.fill.background()
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.color.rgb = BORDER
    return c


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(WIDTH), Inches(HEIGHT))
    base.fill.solid()
    base.fill.fore_color.rgb = RGBColor(15, 18, 27)
    base.line.fill.background()

    vignette = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(-0.7), Inches(15.8), Inches(9.2))
    vignette.fill.solid()
    vignette.fill.fore_color.rgb = RGBColor(24, 27, 38)
    vignette.fill.transparency = 0.58
    vignette.line.fill.background()

    panel_shadow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.18), Inches(1.22), Inches(10.95), Inches(4.7))
    panel_shadow.fill.solid()
    panel_shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    panel_shadow.fill.transparency = 0.58
    panel_shadow.line.fill.background()

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.12), Inches(1.16), Inches(10.95), Inches(4.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(40, 41, 49)
    panel.fill.transparency = 0.14
    panel.line.color.rgb = RGBColor(70, 74, 88)

    style_text(
        s.shapes.add_textbox(Inches(2.0), Inches(1.58), Inches(9.2), Inches(0.95)),
        "Teacher Copilot",
        size=60,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=FONT_SCRIPT,
    )
    style_text(
        s.shapes.add_textbox(Inches(1.8), Inches(2.72), Inches(9.7), Inches(0.6)),
        "AI-Powered Teacher Productivity & Student Learning Insights",
        size=24,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=FONT_SCRIPT,
    )

    map_shadow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.05), Inches(3.72), Inches(5.4), Inches(2.0))
    map_shadow.fill.solid()
    map_shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    map_shadow.fill.transparency = 0.78
    map_shadow.line.fill.background()

    map_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(3.68), Inches(5.4), Inches(2.0))
    map_card.fill.solid()
    map_card.fill.fore_color.rgb = RGBColor(245, 245, 245)
    map_card.line.color.rgb = RGBColor(220, 220, 220)

    style_text(
        s.shapes.add_textbox(Inches(4.28), Inches(4.53), Inches(4.82), Inches(0.5)),
        "Pan-African rollout",
        size=28,
        bold=False,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        italic=True,
        font=FONT_SCRIPT,
    )

    for x, y in [(4.55, 4.15), (5.5, 4.45), (6.42, 4.05), (7.2, 4.42), (8.12, 4.25), (8.82, 4.62)]:
        dot_shadow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.02), Inches(y + 0.02), Inches(0.23), Inches(0.23))
        dot_shadow.fill.solid()
        dot_shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
        dot_shadow.fill.transparency = 0.65
        dot_shadow.line.fill.background()

        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.23), Inches(0.23))
        d.fill.solid()
        d.fill.fore_color.rgb = RGBColor(217, 119, 6)
        d.line.fill.background()

    style_text(
        s.shapes.add_textbox(Inches(2.0), Inches(7.03), Inches(9.3), Inches(0.34)),
        "Empowering teachers and transforming learning outcomes",
        size=24,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=FONT_SCRIPT,
    )


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "The Problem")
    card(s, 0.55, 1.28, 6.25, 5.4)
    style_bullets(
        s.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(5.6), Inches(4.6)),
        [
            "Teachers overloaded with manual grading",
            "Delayed feedback and generic responses",
            "Limited insights for personalized learning",
            "Students unaware of learning gaps",
        ],
        size=24,
        spacing=14,
    )
    add_photo(s, "problem", 7.0, 1.3, 5.75, 5.35)
    add_tagline(s, "Empowering Teachers, Transforming Learning")


def slide_mission(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, soft=True)
    add_header(s, "Our Mission & Vision")
    card(s, 0.55, 1.35, 6.25, 5.25, color=CARD_BG)

    box = s.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.7), Inches(4.4))
    tf = box.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "Mission: Empower teachers to transform learning with AI"
    p1.font.size = Pt(25)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT
    p1.font.name = FONT_HEAD

    p2 = tf.add_paragraph()
    p2.text = "Vision: A future where every student receives personalized education"
    p2.space_before = Pt(30)
    p2.font.size = Pt(25)
    p2.font.bold = True
    p2.font.color.rgb = TEXT
    p2.font.name = FONT_HEAD

    add_photo(s, "mission", 7.0, 1.48, 5.75, 4.95)
    add_tagline(s, "Empowering Teachers, Transforming Learning", color=PRIMARY)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Our Solution:", right_text="Teacher Copilot")

    style_text(s.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(4.6), Inches(0.4)), "Current Workflow", size=17, bold=True, color=MUTED, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    style_text(s.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(4.6), Inches(0.4)), "Optimized Workflow", size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_HEAD)

    current = ["Manual Grading", "Basic Feedback", "Slow Reports"]
    future = ["AI Grading", "Instant Feedback", "Analytics & Alerts"]

    y = 1.85
    for idx in range(3):
        l = card(s, 0.8, y, 4.9, 1.0, color=RGBColor(244, 244, 245))
        l.line.color.rgb = BORDER
        style_text(s.shapes.add_textbox(Inches(1.05), Inches(y + 0.31), Inches(4.4), Inches(0.35)), current[idx], size=18, bold=True, color=TEXT, align=PP_ALIGN.CENTER, font=FONT_HEAD)

        r = card(s, 7.0, y, 5.3, 1.0, color=ACCENT_SOFT)
        r.line.color.rgb = RGBColor(191, 219, 254)
        style_text(s.shapes.add_textbox(Inches(7.3), Inches(y + 0.31), Inches(4.7), Inches(0.35)), future[idx], size=18, bold=True, color=RGBColor(30, 64, 175), align=PP_ALIGN.CENTER, font=FONT_HEAD)

        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.82), Inches(y + 0.24), Inches(1.0), Inches(0.52))
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(161, 161, 170)
        arr.line.fill.background()
        y += 1.43


def slide_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Market Opportunity")

    card(s, 0.6, 1.45, 6.1, 5.0)
    tiers = [
        ("TAM", "2.5M Schools\n$600M+", RGBColor(39, 39, 42), 1.3, 2.0, 4.7, 1.2),
        ("SAM", "Targetable 32M Schools", RGBColor(82, 82, 91), 1.6, 3.28, 4.1, 1.1),
        ("SOM", "Pilot 5 Schools\n$20M", RGBColor(113, 113, 122), 1.9, 4.43, 3.5, 1.0),
    ]
    for label, txt, color, l, t, w, h in tiers:
        shp = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        style_text(s.shapes.add_textbox(Inches(l + 0.2), Inches(t + 0.17), Inches(w - 0.4), Inches(h - 0.25)), f"{label}\n{txt}", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_HEAD)

    add_photo(s, "market", 7.05, 1.53, 5.65, 4.95)


def slide_business(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Business Model")

    items = [
        (0.9, CARD_BG, "B2B (Schools)\n$6,000 per school/year"),
        (6.9, CARD_BG, "B2G (Government)\n$600,000 per country/year"),
    ]
    for idx, (x, color, txt) in enumerate(items):
        c = card(s, x, 1.8, 5.5, 3.35, color=color)
        c.line.color.rgb = BORDER
        style_text(s.shapes.add_textbox(Inches(x + 0.25), Inches(2.23), Inches(5.0), Inches(1.55)), txt, size=25, bold=True, color=TEXT, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.72), Inches(5.5), Inches(0.42))
        strip.fill.solid()
        strip.fill.fore_color.rgb = ACCENT if idx == 0 else RGBColor(99, 102, 241)
        strip.line.fill.background()


def slide_value(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Value Proposition")

    card(s, 0.6, 1.5, 6.3, 5.0, color=ACCENT_SOFT)
    style_bullets(
        s.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.8), Inches(4.5)),
        [
            "For Teachers: Save time and actionable insights",
            "For Students: Personalized learning support",
            "For Schools & Governments: Scalable cost efficiency",
        ],
        size=23,
        spacing=16,
    )
    add_photo(s, "value", 7.2, 1.65, 5.4, 4.85)


def slide_traction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Our Traction")

    card(s, 0.6, 1.5, 6.2, 5.0)
    style_bullets(
        s.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.7), Inches(4.3)),
        [
            "5 schools launched in Rwanda & Ghana",
            "Pipeline: 10 schools in Senegal & Kenya",
            "Positive pilot results with high teacher satisfaction",
        ],
        size=23,
        spacing=16,
    )
    add_photo(s, "traction", 6.95, 1.65, 5.7, 4.85)


def slide_revenue(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Revenue Projections")

    rows = [("Year 1", "$5.4M", RGBColor(63, 63, 70)), ("Year 2", "$14.4M", PRIMARY), ("Year 3", "$248M+", ACCENT)]
    y = 1.8
    for year, amt, color in rows:
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(6.2), Inches(1.1))
        r.fill.solid()
        r.fill.fore_color.rgb = color
        r.line.fill.background()
        style_text(s.shapes.add_textbox(Inches(0.9), Inches(y + 0.3), Inches(2.1), Inches(0.45)), year, size=24, bold=True, color=WHITE, font=FONT_HEAD)
        style_text(s.shapes.add_textbox(Inches(2.35), Inches(y + 0.29), Inches(4.0), Inches(0.45)), amt, size=31, bold=True, color=WHITE, font=FONT_HEAD)
        y += 1.19

    b_x = 8.1
    for i, h in enumerate([1.2, 1.95, 2.8]):
        c = [RGBColor(161, 161, 170), RGBColor(113, 113, 122), ACCENT][i]
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(b_x + 0.78 * i), Inches(6.0 - h), Inches(0.56), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()

    up = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(7.15), Inches(3.2), Inches(4.9), Inches(2.95))
    up.fill.solid()
    up.fill.fore_color.rgb = RGBColor(59, 130, 246)
    up.fill.transparency = 0.2
    up.line.fill.background()


def slide_gtm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Go-To-Market Strategy")

    labels = [("Year 1", "$5.4M", 1.0), ("Year 2", "$14.4M", 4.3), ("Year 3", "$248M+", 7.9)]
    for year, amt, x in labels:
        style_text(s.shapes.add_textbox(Inches(x), Inches(5.45), Inches(2.2), Inches(0.36)), year, size=20, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        style_text(s.shapes.add_textbox(Inches(x), Inches(3.24), Inches(2.5), Inches(0.8)), amt, size=31, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)

    bars = [
        (1.2, 0.7, RGBColor(113, 113, 122)),
        (2.0, 0.9, RGBColor(82, 82, 91)),
        (4.8, 1.6, RGBColor(82, 82, 91)),
        (5.58, 1.6, PRIMARY),
        (8.2, 2.15, RGBColor(96, 165, 250)),
        (9.0, 2.75, RGBColor(59, 130, 246)),
        (9.78, 3.2, ACCENT),
    ]
    for x, h, col in bars:
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.3 - h), Inches(0.54), Inches(h))
        b.fill.solid()
        b.fill.fore_color.rgb = col
        b.line.fill.background()

    growth = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(0.8), Inches(2.0), Inches(11.45), Inches(3.2))
    growth.fill.solid()
    growth.fill.fore_color.rgb = RGBColor(59, 130, 246)
    growth.fill.transparency = 0.22
    growth.line.fill.background()


def slide_technology(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Technology Overview")
    card(s, 0.6, 1.5, 6.2, 5.0)

    items = ["AI Grading Engine", "Analytics Dashboard", "Cloud-Based Platform", "LMS Integrations"]
    y = 1.86
    for item in items:
        icon = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.36), Inches(0.36))
        icon.fill.solid()
        icon.fill.fore_color.rgb = ACCENT
        icon.line.fill.background()
        style_text(s.shapes.add_textbox(Inches(1.35), Inches(y - 0.02), Inches(5.2), Inches(0.42)), item, size=23, bold=True, color=TEXT, font=FONT_HEAD)
        y += 1.08

    add_photo(s, "technology", 7.0, 1.65, 5.65, 4.85)


def slide_impact(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Impact & Metrics")
    card(s, 0.6, 1.5, 6.2, 5.0)

    style_bullets(
        s.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.7), Inches(4.3)),
        [
            "Teacher hours saved",
            "Student engagement +40%",
            "Faster feedback: 75% reduction",
            "Improved outcomes across pilot schools",
        ],
        size=23,
        spacing=16,
    )

    data = CategoryChartData()
    data.categories = ["Teacher Time", "Engagement", "Other"]
    data.add_series("Impact", (50, 40, 10))
    chart = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(7.35), Inches(1.7), Inches(5.1), Inches(4.9), data).chart
    chart.has_legend = False
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.number_format = "0%"
    chart.plots[0].data_labels.position = 2
    fills = [ACCENT, RGBColor(96, 165, 250), RGBColor(191, 219, 254)]
    for i, p in enumerate(chart.series[0].points):
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = fills[i]


def slide_funding(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Funding Ask")
    card(s, 0.6, 1.5, 5.1, 5.0)

    style_bullets(
        s.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(4.6), Inches(4.4)),
        [
            "Raising $3M",
            "Product development",
            "Sales & marketing",
            "School onboarding",
            "Operations",
        ],
        size=22,
        spacing=14,
    )

    d = CategoryChartData()
    d.categories = ["Product", "Sales", "Onboarding", "Operations"]
    d.add_series("Allocation", (40, 30, 20, 10))
    pie = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(5.65), Inches(2.0), Inches(3.9), Inches(3.9), d).chart
    pie.has_legend = False
    pie.plots[0].has_data_labels = True
    pie.plots[0].data_labels.number_format = "0%"
    pie.plots[0].data_labels.position = 2
    colors = [PRIMARY, ACCENT, RGBColor(96, 165, 250), RGBColor(191, 219, 254)]
    for i, p in enumerate(pie.series[0].points):
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = colors[i]

    add_photo(s, "market", 9.65, 1.7, 2.95, 4.8)


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Our Team")

    people = [
        ("team_1", "Founder & CEO"),
        ("team_2", "CTO"),
        ("team_3", "Head of Partnerships"),
        ("team_4", "Education Advisors"),
    ]
    x = 0.8
    for key, role in people:
        card(s, x, 1.6, 2.95, 5.1, color=RGBColor(250, 251, 252))
        add_photo(s, key, x + 0.18, 1.86, 2.6, 3.4, rounded=False)
        style_text(
            s.shapes.add_textbox(Inches(x + 0.1), Inches(5.5), Inches(2.75), Inches(0.9)),
            role,
            size=16,
            bold=True,
            align=PP_ALIGN.CENTER,
            font=FONT_HEAD,
        )
        x += 3.12


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, soft=True)
    add_header(s, "Join Us to", right_text="Transform African Education")

    card(s, 0.65, 1.55, 6.2, 4.6, color=CARD_BG)
    style_text(
        s.shapes.add_textbox(Inches(1.05), Inches(2.15), Inches(5.35), Inches(2.3)),
        "Partner with us\nTo empower teachers &\nstudents across Africa",
        size=30,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
        italic=True,
        font=FONT_HEAD,
    )

    add_photo(s, "closing", 7.0, 1.8, 5.7, 3.95)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.45), Inches(WIDTH), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    style_text(
        s.shapes.add_textbox(Inches(0.7), Inches(6.72), Inches(12.0), Inches(0.5)),
        "Contact Us: email@example.com  |  www.teachercopilot.com",
        size=19,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=FONT_HEAD,
    )


def build():
    fetch_images()
    prs = Presentation()
    prs.slide_width = Inches(WIDTH)
    prs.slide_height = Inches(HEIGHT)

    slide_cover(prs)
    slide_problem(prs)
    slide_mission(prs)
    slide_solution(prs)
    slide_market(prs)
    slide_business(prs)
    slide_value(prs)
    slide_traction(prs)
    slide_revenue(prs)
    slide_gtm(prs)
    slide_technology(prs)
    slide_impact(prs)
    slide_funding(prs)
    slide_team(prs)
    slide_closing(prs)

    candidates = [
        "Teacher_Copilot_Pitch_Deck.pptx",
        "Teacher_Copilot_Pitch_Deck_Enhanced.pptx",
        f"Teacher_Copilot_Pitch_Deck_shadcn_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
    ]
    for name in candidates:
        try:
            prs.save(name)
            return
        except PermissionError:
            continue
    raise PermissionError("Could not save deck. Close open .pptx files and rerun.")


if __name__ == "__main__":
    build()
